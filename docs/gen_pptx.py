#!/usr/bin/env python3
"""Generate Spring Boot + Spring Cloud training PPT."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
SPRING_GREEN = RGBColor(0x6D, 0xB3, 0x3F)
SPRING_DARK = RGBColor(0x34, 0x30, 0x2D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
RED = RGBColor(0xE7, 0x4C, 0x3C)
GRAY = RGBColor(0x66, 0x66, 0x66)
DARK_BG = RGBColor(0x1B, 0x1B, 0x1B)
LIGHT_GREEN_BG = RGBColor(0xF0, 0xFD, 0xF4)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 color=SPRING_DARK, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_code_box(slide, left, top, width, height, code_text, font_size=11):
    """Add a dark code block."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x27, 0x28, 0x22)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    for i, line in enumerate(code_text.strip().split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = 'Courier New'
        p.font.color.rgb = RGBColor(0xE6, 0xDB, 0x74)
        p.space_after = Pt(2)
    return tf

def add_table(slide, left, top, width, height, rows, col_widths=None):
    """Add a styled table. rows is list of lists."""
    n_rows = len(rows)
    n_cols = len(rows[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for r, row in enumerate(rows):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.name = 'Microsoft YaHei'
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = SPRING_DARK
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = SPRING_GREEN
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 0 else LIGHT_BG
    return table

def add_chapter_slide(title, chap_num, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, DARK_BG)
    # Large number
    add_text_box(slide, 0.8, 1.2, 4, 2, chap_num, font_size=96, bold=True,
                 color=RGBColor(0x6D, 0xB3, 0x3F), alignment=PP_ALIGN.LEFT)
    # Title
    add_text_box(slide, 0.8, 3.5, 10, 1.5, title, font_size=40, bold=True,
                 color=WHITE, alignment=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, 0.8, 4.8, 10, 0.8, subtitle, font_size=20,
                     color=RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.LEFT)
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.2), Inches(1.5), Inches(0.06))
    line.fill.solid(); line.fill.fore_color.rgb = SPRING_GREEN; line.line.fill.background()
    return slide

def add_compare_slide(title, bad_title, bad_items, good_title, good_items):
    """Two-column comparison slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, 0.8, 0.4, 11, 0.8, title, font_size=30, bold=True, color=SPRING_DARK)
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(1.2), Inches(0.05))
    line.fill.solid(); line.fill.fore_color.rgb = SPRING_GREEN; line.line.fill.background()
    # Bad column
    bad = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2))
    bad.fill.solid(); bad.fill.fore_color.rgb = RGBColor(0xFF, 0xF5, 0xF5)
    bad.line.color.rgb = RGBColor(0xFE, 0xCA, 0xCA); bad.line.width = Pt(2)
    add_text_box(slide, 1.1, 1.8, 5, 0.5, bad_title, font_size=22, bold=True, color=RED)
    for i, item in enumerate(bad_items):
        add_text_box(slide, 1.1, 2.5 + i * 0.6, 5, 0.5, f"• {item}", font_size=14, color=SPRING_DARK)
    # Good column
    good = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.2))
    good.fill.solid(); good.fill.fore_color.rgb = LIGHT_GREEN_BG
    good.line.color.rgb = RGBColor(0xBB, 0xF7, 0xD0); good.line.width = Pt(2)
    add_text_box(slide, 7.3, 1.8, 5, 0.5, good_title, font_size=22, bold=True, color=SPRING_GREEN)
    for i, item in enumerate(good_items):
        add_text_box(slide, 7.3, 2.5 + i * 0.6, 5, 0.5, f"• {item}", font_size=14, color=SPRING_DARK)
    return slide

def add_card_slide(title, cards):
    """cards = [(emoji, title, desc), ...]"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, 0.8, 0.4, 11, 0.8, title, font_size=30, bold=True, color=SPRING_DARK)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(1.2), Inches(0.05))
    line.fill.solid(); line.fill.fore_color.rgb = SPRING_GREEN; line.line.fill.background()
    n = len(cards)
    card_w = min(3.5, 11.5 / n - 0.3)
    start_x = (12.5 - (card_w + 0.3) * n) / 2
    for i, (emoji, ct, desc) in enumerate(cards):
        x = start_x + i * (card_w + 0.3)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(2.0), Inches(card_w), Inches(4.0))
        shape.fill.solid(); shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0); shape.line.width = Pt(1)
        add_text_box(slide, x + 0.2, 2.3, card_w - 0.4, 1, emoji, font_size=36, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.2, 3.3, card_w - 0.4, 0.6, ct, font_size=18, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.2, 3.9, card_w - 0.4, 1.5, desc, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)
    return slide

# ========== TITLE SLIDE ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_text_box(slide, 1, 1.5, 11, 1.5, "Spring Boot + Spring Cloud", font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.2, 11, 0.8, "产品团队技术培训", font_size=28, color=RGBColor(0xCC,0xCC,0xCC), alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.3, 11, 0.6, "Spring Boot 3.x  ·  MyBatis Plus  ·  Spring Cloud", font_size=16, color=SPRING_GREEN, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.5, 11, 0.5, "~2h35min  ·  上半场实操 + 下半场理论", font_size=14, color=RGBColor(0x88,0x88,0x88), alignment=PP_ALIGN.CENTER)

# ========== AGENDA ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "📋 今天讲什么", font_size=30, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(1.2), Inches(0.05))
line.fill.solid(); line.fill.fore_color.rgb = SPRING_GREEN; line.line.fill.background()
# Left card
left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5))
left.fill.solid(); left.fill.fore_color.rgb = LIGHT_GREEN_BG
left.line.color.rgb = RGBColor(0xBB, 0xF7, 0xD0); left.line.width = Pt(2)
add_text_box(slide, 1.2, 2.1, 5, 0.5, "🛠️ 上半场：Spring Boot 实操", font_size=22, bold=True, color=SPRING_GREEN)
add_text_box(slide, 1.2, 2.9, 5, 2, "演进之路 → 核心机制 → 项目结构\n→ MyBatis Plus → 配置管理\n→ 统一响应 → AI 实战 → Demo", font_size=14, color=SPRING_DARK)
add_text_box(slide, 1.2, 5.2, 5, 0.4, "~100 min", font_size=16, color=GRAY)
# Right card
right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.5))
right.fill.solid(); right.fill.fore_color.rgb = LIGHT_GREEN_BG
right.line.color.rgb = RGBColor(0xBB, 0xF7, 0xD0); right.line.width = Pt(2)
add_text_box(slide, 7.4, 2.1, 5, 0.5, "☁️ 下半场：Spring Cloud 理论", font_size=22, bold=True, color=SPRING_GREEN)
add_text_box(slide, 7.4, 2.9, 5, 2, "为什么微服务 → 六大组件\n→ 产品视角 → 总结", font_size=14, color=SPRING_DARK)
add_text_box(slide, 7.4, 5.2, 5, 0.4, "~58 min", font_size=16, color=GRAY)

# ========== CH1: 演进之路 ==========
add_chapter_slide("从 Spring MVC 到 Spring Boot", "01", "演进之路 · ~15min")

add_compare_slide("传统 Spring 开发有多痛？",
    "😫 四大痛点", [
        "配置地狱：web.xml + applicationContext.xml + spring-mvc.xml",
        "依赖管理：6~8 个依赖，版本自己对齐",
        "部署麻烦：装 Tomcat → 打 WAR → 丢 webapps → 重启",
        "启动慢：搭环境半天起步",
    ],
    "💀 还没写代码…", [
        "1. web.xml (至少30行XML)",
        "2. applicationContext.xml (至少20行XML)",
        "3. spring-mvc.xml (至少30行XML)",
        "4. 配置数据源、事务管理器...",
        "5. 安装 Tomcat → 6. 打WAR → 7. 部署",
        "→ 一行业务代码还没写，100+ 行 XML 先写出来",
    ])

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "💡 Spring Boot 的答案：约定大于配置", font_size=30, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(1.2), Inches(0.05))
line.fill.solid(); line.fill.fore_color.rgb = SPRING_GREEN; line.line.fill.background()
add_table(slide, 0.8, 1.8, 11.5, 3.5, [
    ["你的动作", "Spring Boot 的约定（理解）", "省掉了什么（价值）"],
    ["引入 spring-webmvc", "\"你要用 Spring MVC\" → 自动配 DispatcherServlet", "web.xml + spring-mvc.xml"],
    ["继承 BaseMapper<Employee>", "\"你要 CRUD 这个表\" → 自动生成 SQL", "百行 XML 映射文件"],
    ["引入 starter-web + 不配端口", "\"默认 8080，内嵌 Tomcat\"", "装 Tomcat、配端口、打 WAR"],
], col_widths=[3, 4.5, 4])
add_text_box(slide, 1, 6.0, 11, 0.8, '把"每次都要写的东西"变成"默认就有的东西"。想改可以改，不改就能用。',
             font_size=18, bold=True, color=SPRING_GREEN, alignment=PP_ALIGN.CENTER)

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "🔧 Spring Initializr — 30秒生成项目", font_size=28, bold=True)
add_text_box(slide, 0.8, 1.5, 11, 0.6, "🌐 网址：start.spring.io", font_size=22, bold=True, color=SPRING_GREEN)
add_table(slide, 0.8, 2.4, 11.5, 3.2, [
    ["步骤", "操作", "填什么"],
    ["1", "打开网页", "start.spring.io"],
    ["2", "构建工具", "Maven"],
    ["3", "语言 & 版本", "Java 21 + Spring Boot 3.3.x"],
    ["4", "项目坐标", "Group: com.example / Artifact: my-first-boot-app"],
    ["5", "加依赖", "Spring Web + Lombok + H2 Database"],
], col_widths=[1.5, 4, 6])
add_text_box(slide, 0.8, 6.0, 11, 0.6, "点「Generate」→ 下载 zip → 解压 → 直接能跑（传统半小时 vs 现在30秒）",
             font_size=18, bold=True, color=SPRING_GREEN, alignment=PP_ALIGN.CENTER)

# ========== CH2: 核心机制 ==========
add_chapter_slide("Spring Boot 核心机制", "02", "自动配置 · Starter · 条件装配 · ~12min")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "🎯 一个注解 = 三件事", font_size=28, bold=True)
add_code_box(slide, 0.8, 1.5, 11.5, 1.5, """@SpringBootApplication  // ← 三合一注解
public class DemoApplication {
    public static void main(String[] args) {
        SpringApplication.run(DemoApplication.class, args);
    }
}""")
add_table(slide, 0.8, 3.5, 11.5, 2.5, [
    ["注解", "作用", "学过？"],
    ["@Configuration", "可以定义 Bean，标记配置类", "✅ 学过的"],
    ["@ComponentScan", "扫描当前包及子包下所有组件", "✅ 学过的"],
    ["@EnableAutoConfiguration", "自动配置——Spring Boot 的灵魂", "🆕 新东西！"],
], col_widths=[4, 4.5, 3])

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "🏠 精装修 vs 毛坯房", font_size=28, bold=True)
add_compare_slide("", "", [], "", [])
# Override with custom content
for sh in list(slide.shapes)[1:]:  # remove auto-added shapes
    pass
add_text_box(slide, 1.0, 1.2, 5, 0.5, "传统 Spring = 毛坯房", font_size=22, bold=True, color=RED)
add_text_box(slide, 1.0, 1.9, 5, 3, "墙自己刷\n电线自己铺\n地板自己装\n\n好处：想怎么装怎么装\n坏处：累", font_size=16, color=SPRING_DARK)
add_text_box(slide, 7.5, 1.2, 5, 0.5, "Spring Boot = 精装修公寓", font_size=22, bold=True, color=SPRING_GREEN)
add_text_box(slide, 7.5, 1.9, 5, 3, "墙刷好了\n空调装了\n热水器接了\n\n拎包入住\n不满意可以改——你的配置优先级更高", font_size=16, color=SPRING_DARK)

# ========== CH3: 项目结构 ==========
add_chapter_slide("Spring Boot 项目结构", "03", "什么没了？什么没变？ · ~5min")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "🔍 找茬——什么没了？", font_size=28, bold=True)
# Bad
bad = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.5))
bad.fill.solid(); bad.fill.fore_color.rgb = RGBColor(0xFF, 0xF5, 0xF5)
bad.line.color.rgb = RGBColor(0xFE, 0xCA, 0xCA); bad.line.width = Pt(2)
add_text_box(slide, 1.1, 1.9, 5, 0.5, "❌ 传统 Spring MVC", font_size=22, bold=True, color=RED)
add_text_box(slide, 1.1, 2.6, 5, 2.5, "webapp/WEB-INF/web.xml\napplicationContext.xml\nspring-mvc.xml\n\n3 个 XML 配置文件", font_size=16, color=SPRING_DARK)
# Good
good = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.6), Inches(5.5), Inches(4.5))
good.fill.solid(); good.fill.fore_color.rgb = LIGHT_GREEN_BG
good.line.color.rgb = RGBColor(0xBB, 0xF7, 0xD0); good.line.width = Pt(2)
add_text_box(slide, 7.3, 1.9, 5, 0.5, "✅ Spring Boot", font_size=22, bold=True, color=SPRING_GREEN)
add_text_box(slide, 7.3, 2.6, 5, 2.5, "DemoApplication.java\napplication.yml\n\n1 个启动类 + 1 个 YAML\n内嵌 Tomcat，java -jar 直接跑", font_size=16, color=SPRING_DARK)
add_text_box(slide, 2, 6.5, 9, 0.5, "分层架构不变：controller/ → service/ → mapper/ → entity/", font_size=18, bold=True, alignment=PP_ALIGN.CENTER)

# ========== CH4: MyBatis Plus ==========
add_chapter_slide("MyBatis Plus", "04", '数据访问层的"自动挡" · ~15min')

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "🎁 继承即拥有 — 空接口 = 全套 CRUD", font_size=28, bold=True)
add_code_box(slide, 0.8, 1.3, 11.5, 2.8, """// Entity: 一个类对应一张表
@Data @TableName("employee")
public class Employee {
    @TableId(type = IdType.AUTO) private Long id;
    private String name; private Integer age; private String email;
}

// Mapper: 继承 BaseMapper，自动获得 CRUD！
@Mapper
public interface EmployeeMapper extends BaseMapper<Employee> {
    // 空接口！但拥有 insert/deleteById/updateById/selectById/
    //            selectList/selectPage... 十几种方法
}""")
add_text_box(slide, 0.8, 4.6, 11, 2, "一张新表 = 一个空接口继承 BaseMapper。完事。\n—— MyBatis Plus 的哲学跟 Spring Boot 一模一样：约定大于配置。",
             font_size=18, bold=True, color=SPRING_GREEN)

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "🔗 LambdaQueryWrapper — 不写 SQL", font_size=28, bold=True)
add_code_box(slide, 0.8, 1.3, 11.5, 2.5, """LambdaQueryWrapper<Employee> wrapper = new LambdaQueryWrapper<>();
wrapper.eq(Employee::getName, "张三")    // WHERE name = '张三'
       .gt(Employee::getAge, 25)         // AND age > 25
       .orderByDesc(Employee::getId);    // ORDER BY id DESC
employeeMapper.selectList(wrapper);

// 分页 —— 一行搞定
Page<Employee> page = new Page<>(1, 10);
employeeMapper.selectPage(page, wrapper);
// 自动返回：records、total、current、pages  """)
add_text_box(slide, 0.8, 4.5, 11, 1.5, "字段写错了？编译器直接报错——不是运行时炸。\n演进路径：JDBC 手写 → MyBatis 半自动 → MyBatis Plus 更自动",
             font_size=16, color=SPRING_DARK)

# ========== CH5: 配置管理 ==========
add_chapter_slide("配置管理", "05", "从 XML 到 YAML · ~5min")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "📉 21 行 XML → 5 行 YAML", font_size=28, bold=True)
add_code_box(slide, 0.8, 1.5, 5, 3.5, """<!-- 传统 XML: 21行 -->
<bean id="dataSource"
      class="...BasicDataSource">
  <property name="driverClassName"
            value="com.mysql.cj..."/>
  <property name="url"
            value="jdbc:mysql://..."/>
  <property name="username"
            value="root"/>
  <property name="password"
            value="123456"/>
</bean>""", font_size=10)
add_code_box(slide, 6.5, 1.5, 5.8, 3.5, """# Spring Boot YAML: 5行
spring:
  datasource:
    url: jdbc:mysql://localhost/demo
    username: root
    password: ${DB_PASSWORD}
    driver-class-name: com.mysql.cj.jdbc.Driver

# 多环境: 一行切换
# spring.profiles.active=prod""", font_size=10)
add_text_box(slide, 0.8, 5.5, 11, 1, "没有尖括号 · 没有 property · 没有 bean 定义\n密码用 ${DB_PASSWORD} 环境变量，永远不要硬编码",
             font_size=16, color=SPRING_DARK)

# ========== CH6: 统一响应与异常 ==========
add_chapter_slide("统一响应与异常处理", "06", "前端同事的救命稻草 · ~10min")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "✅ 统一用 Result<T> 包装", font_size=28, bold=True)
add_code_box(slide, 0.8, 1.3, 11.5, 1.2, """{"code": 200, "message": "操作成功", "data": { "id": 1, "name": "张三" }}""", font_size=16)
add_text_box(slide, 0.8, 2.8, 11, 0.5, "三个字段：code（状态码）+ message（提示）+ data（真正的数据，泛型）", font_size=16)
add_code_box(slide, 0.8, 3.5, 5.5, 2.5, """// Controller 统一返回 Result<T>
@GetMapping("/{id}")
public Result<Employee> getById(@PathVariable Long id) {
    return Result.success(employeeService.getById(id));
}

@DeleteMapping("/{id}")
public Result<Void> delete(@PathVariable Long id) {
    employeeService.delete(id);
    return Result.success();
}""", font_size=10)
add_code_box(slide, 7.0, 3.5, 5.3, 2.5, """// 全局异常处理 — 一个类搞定
@RestControllerAdvice
public class GlobalExceptionHandler {
    @ExceptionHandler(BusinessException.class)
    public Result<?> handle(BusinessException e) {
        return Result.error(e.getCode(), e.getMessage());
    }
}
// Controller 里不需要写任何 try-catch""", font_size=10)
add_text_box(slide, 0.8, 6.3, 11, 0.6, "审 AI 代码三件事：返回 Result<T>？抛 BusinessException？有 GlobalExceptionHandler？",
             font_size=16, bold=True, color=SPRING_GREEN)

# ========== CH7: AI 实战 ==========
add_chapter_slide("AI 辅助开发实战指南", "⭐ 07", "Prompt 模板 · 15条审查清单 · ~20min")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "📝 给 AI 写 Prompt 的模板", font_size=28, bold=True)
add_code_box(slide, 0.8, 1.3, 11.5, 3.8, """使用 Spring Boot 3.x + MyBatis Plus 开发【功能描述】。

要求：
- 分层架构：Controller / Service / ServiceImpl / Mapper / Entity
- 统一返回 Result<T>
- 增删改加 @Transactional
- 用 LambdaQueryWrapper，不要字符串拼 SQL
- 日期用 LocalDateTime，金额用 BigDecimal
- 异常用 BusinessException，不要空 catch
- 包名用 jakarta，不要 javax
- 不要生成 XML 配置，全用注解和 YAML""")
add_text_box(slide, 0.8, 5.6, 11, 1.2, "每一个要求前面几章都讲过了。你每多说一句要求，就少踩一个坑。\n关键：AI 不会理解你项目的上下文——你说了才有。",
             font_size=16, color=SPRING_DARK)

# 15条审查清单 - 3 slides
for batch, items in enumerate([
    ["① 文件在正确包下？", "② Controller 有 @RestController？", "③ Service 接口+实现分离？",
     "④ Mapper 继承 BaseMapper？", "⑤ @TableName 对应表名？"],
    ["⑥ 返回值是 Result<T>？", "⑦ 增删改有 @Transactional？", "⑧ 密码硬编码了？",
     "⑨ 用了 LambdaQueryWrapper？", "⑩ catch 块是空的？"],
    ["⑪ 分页用了 Page<T>？", "⑫ 日期用 LocalDateTime？", "⑬ 参数有 @Valid？",
     "⑭ 日志用 @Slf4j？", "⑮ import 有 javax？(应为 jakarta)"],
]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, 0.8, 0.4, 11, 0.8, f"🔍 15 条审查清单（{batch+1}/3）", font_size=28, bold=True)
    for i, item in enumerate(items):
        y = 1.8 + i * 1.0
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(y), Inches(10), Inches(0.7))
        shape.fill.solid(); shape.fill.fore_color.rgb = LIGHT_BG
        shape.line.color.rgb = SPRING_GREEN; shape.line.width = Pt(1)
        add_text_box(slide, 1.8, y + 0.1, 9.5, 0.5, item, font_size=18, color=SPRING_DARK)
    if batch == 2:
        add_text_box(slide, 1.5, 6.8, 10, 0.5, "这 15 条，不需要会写代码就能查 —— AI 写代码，你质检",
                     font_size=18, bold=True, color=SPRING_GREEN, alignment=PP_ALIGN.CENTER)

# AI 3 大错误
add_card_slide("🐛 AI 最常犯的 3 个错误", [
    ("🔄", "循环里查数据库", "for 循环里 selectById\n→ 应该 selectBatchIds"),
    ("💥", "忘了 @Transactional", "转账扣钱成功加钱失败\n→ 钱凭空消失"),
    ("🕳️", "空 catch", "catch(Exception e){}\n→ 定时炸弹"),
])

# ========== CH8: DEMO ==========
add_chapter_slide("Demo 现场演示", "08", "员工管理系统 · ~15min")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "🚀 启动 & 测试", font_size=28, bold=True)
add_code_box(slide, 0.8, 1.3, 5.5, 2.8, """# 启动项目
cd demo
mvn spring-boot:run
# Started DemoApplication in 0.958s

# H2 控制台
# http://localhost:8080/h2-console
# JDBC URL: jdbc:h2:mem:demo""", font_size=10)
add_code_box(slide, 7.0, 1.3, 5.3, 2.8, """# 查所有部门
curl http://localhost:8080/api/departments

# 分页查员工
curl "localhost:8080/api/employees?current=1&size=5"

# 新增员工
curl -X POST http://localhost:8080/api/employees \\
  -H "Content-Type: application/json" \\
  -d '{"name":"新员工","age":26,"salary":15000}'""", font_size=10)
add_text_box(slide, 0.8, 4.6, 11, 2, "不到 1 秒启动。完整 Web 应用 — 数据库 + HTTP + JSON 全在内存里。\nH2 内存数据库，应用一停数据就没了——非常适合演示。",
             font_size=16, color=SPRING_DARK)

# ========== CH9: 为什么微服务 ==========
add_chapter_slide("从单体到微服务", "09", "为什么需要 Spring Cloud · ~12min")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "🏢 一个电商公司的故事", font_size=28, bold=True)
add_text_box(slide, 0.8, 1.5, 11, 0.5, "创业第一天：一个 shop.jar — 订单、用户、商品全在里面。真爽。这就是单体架构。", font_size=18)
add_text_box(slide, 0.8, 2.3, 11, 0.5, "一年后 200 人：", font_size=22, bold=True, color=RED)
items_text = [
    "😤 \"昨天写的代码被谁覆盖了？\" 合并冲突 50 个文件",
    "😤 \"订单扩容 → 整个系统部署 10 份。\" 4G×10=40G 内存，实际上只有订单需要扩容",
    "😤 \"改一行代码 → 全站重新部署。\" 换灯泡拉整栋楼电闸",
    "😤 \"用户模块挂了 → 全站崩溃。\" 一个房间短路，整栋楼停电",
]
for i, t in enumerate(items_text):
    add_text_box(slide, 1.2, 3.0 + i * 0.7, 10, 0.6, t, font_size=16, color=SPRING_DARK)

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "⚠️ 拆完之后的新问题", font_size=28, bold=True)
add_text_box(slide, 0.8, 1.9, 11, 0.6, "微服务的答案：把一个超级工程拆成几十个小工程，每个独立开发、独立部署、独立扩容。", font_size=18, bold=True)
items_q = [
    "❓ 订单服务要调用户服务 —— 怎么找到它在哪台机器？",
    "❓ 前端要调 30 个服务 —— 难道记住 30 个地址？",
    "❓ 用户服务挂了 —— 怎么不让它把订单服务也拖死？",
    "❓ 一个请求跨 5 个服务 —— 慢了怎么排查？",
]
for i, t in enumerate(items_q):
    add_text_box(slide, 1.2, 3.0 + i * 0.8, 10, 0.6, t, font_size=18, color=SPRING_DARK)
add_text_box(slide, 0.8, 6.3, 11, 0.6, "Spring Boot 让一个服务跑起来 → Spring Cloud 让一堆服务能协作",
             font_size=22, bold=True, color=SPRING_GREEN, alignment=PP_ALIGN.CENTER)

# ========== CH10: 六大组件 ==========
add_chapter_slide("Spring Cloud 六大组件", "10", "拆完之后的六道保险 · ~28min")

components = [
    ("📒", "Nacos 注册中心", "公司通讯录", "每个服务启动去登记，调用方去查询。\n你换手机号—更新通讯录就行，不用群发500人。"),
    ("🏢", "Gateway 网关", "写字楼前台", "前端只调 Gateway 一个入口。\n统一鉴权：没 token 的门都进不来。"),
    ("☁️", "Nacos Config 配置中心", "iCloud 同步", "一处修改，全部推送，不重启生效。\n换新手机→iCloud→全部自动同步。"),
    ("📞", "Feign 服务调用", "快捷拨号", "写接口+注解=自动发 HTTP。\n设一个\"老婆\"快捷拨号，说名字就拨。"),
    ("⚡", "Sentinel 熔断降级", "保险丝", "下游挂了→熔断切断+降级兜底。\n电饭煲短路→保险丝断→只断厨房，客厅还亮。"),
    ("📦", "Sleuth/Zipkin 链路追踪", "快递单号", "每个请求一个 TraceID 串联全链路。\n快递慢了？看物流轨迹—卡分拨中心5h→破案。"),
]
for i, (emoji, name, metaphor, desc) in enumerate(components):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    # Left: emoji + name
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(5.5), Inches(5.0))
    shape.fill.solid(); shape.fill.fore_color.rgb = LIGHT_GREEN_BG
    shape.line.color.rgb = RGBColor(0xBB, 0xF7, 0xD0); shape.line.width = Pt(2)
    add_text_box(slide, 1.2, 1.8, 5, 1, emoji, font_size=48, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1.2, 3.0, 5, 0.6, name, font_size=24, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1.2, 3.7, 5, 0.5, f'「{metaphor}」', font_size=20, color=SPRING_GREEN, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1.2, 4.5, 5, 1, desc, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)
    # Right: problem/solution
    add_text_box(slide, 7.0, 1.5, 5.5, 0.5, "❓ 问题", font_size=20, bold=True, color=RED)
    problem = {
        "📒": "订单服务怎么知道用户服务在哪台机器？",
        "🏢": "前端记30个地址？每个服务都做鉴权？",
        "☁️": "数据库地址变了→改10份配置重启10个服务？",
        "📞": "手写HTTP：拼URL、发请求、解析响应？",
        "⚡": "用户服务挂了→订单不停重试→订单也挂→全站崩溃？",
        "📦": "请求跨5个服务，慢了——哪一步的问题？",
    }
    solution = {
        "📒": "Nacos 登记+查询。自动感知服务上下线。",
        "🏢": "单一入口+统一鉴权。前端只记一个地址。",
        "☁️": "配置中心统一管理。一处改，全部自动生效。",
        "📞": "声明式调用。像调本地方法一样调远程。",
        "⚡": "熔断（切断）+降级（兜底）。自动恢复。",
        "📦": "全局 TraceID。调用链每个节点耗时一目了然。",
    }
    add_text_box(slide, 7.2, 2.2, 5.2, 1.2, problem.get(emoji, ""), font_size=14, color=SPRING_DARK)
    add_text_box(slide, 7.0, 3.8, 5.5, 0.5, "✅ 解法", font_size=20, bold=True, color=SPRING_GREEN)
    add_text_box(slide, 7.2, 4.5, 5.2, 1.2, solution.get(emoji, ""), font_size=14, color=SPRING_DARK)

# Components summary grid
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "🧩 六大组件速查", font_size=28, bold=True)
for i, (emoji, name, metaphor, _) in enumerate(components):
    row = i // 3
    col = i % 3
    x = 0.8 + col * 4.1
    y = 1.6 + row * 2.8
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.8), Inches(2.4))
    shape.fill.solid(); shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0); shape.line.width = Pt(1)
    add_text_box(slide, x + 0.3, y + 0.2, 3.2, 0.7, f"{emoji} {name}", font_size=18, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.3, y + 1.0, 3.2, 0.5, f"「{metaphor}」", font_size=14, color=SPRING_GREEN, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 0.8, 7.0, 11, 0.4, "六个组件，六个生活场景。你不需要会配置，但开会时不会一脸懵。",
             font_size=16, bold=True, color=SPRING_GREEN, alignment=PP_ALIGN.CENTER)

# ========== CH11: 产品视角 ==========
add_chapter_slide("产品视角", "11", "微服务不是免费的午餐 · ~8min")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "❓ 拆之前，先问三个问题", font_size=28, bold=True)
questions = [
    ("① 有独立的业务边界吗？", "不是\"功能多\"就该拆。订单和订单详情是同一个边界。"),
    ("② 需要独立扩容吗？", "双十一只有订单要加机器，用户不用——才值得拆。"),
    ("③ 有独立团队维护吗？", "一个团队维护3个微服务 = 分布式单体，更复杂还没拿到好处。"),
]
for i, (q, detail) in enumerate(questions):
    y = 1.5 + i * 1.7
    add_text_box(slide, 1.2, y, 10, 0.6, q, font_size=24, bold=True, color=SPRING_DARK)
    add_text_box(slide, 1.6, y + 0.7, 10, 0.5, detail, font_size=16, color=GRAY)
add_text_box(slide, 1.2, 6.3, 10, 0.6, "每拆一个服务，运维成本就多一份。拆之前，先问这三个问题。",
             font_size=20, bold=True, color=SPRING_GREEN, alignment=PP_ALIGN.CENTER)

# ========== CH12: 总结 ==========
add_chapter_slide("课程总结 + Q&A", "12", "~10min")

add_card_slide("🚗 你最该带走的三件事", [
    ("🚗", "Spring Boot = 自动挡", "学过的 Spring 全部有效\n只是不需要 XML 配置了"),
    ("📋", "15 条审查清单", "你不会写，但你会查\nAI 写代码，你质检"),
    ("🏗️", "Cloud 六道保险", "拆之前先三问\n不是免费的午餐"),
])

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "📝 课后资料", font_size=28, bold=True, alignment=PP_ALIGN.CENTER)
items = [
    "📂 Demo 源码：demo/",
    "▶️ 启动命令：cd demo && mvn spring-boot:run",
    "📄 大纲 & 逐字稿：docs/",
    "🌐 Spring Initializr：start.spring.io",
    "📋 AI 审查清单 —— 见大纲第7章",
]
for i, t in enumerate(items):
    add_text_box(slide, 3, 1.8 + i * 0.9, 7, 0.6, t, font_size=20, color=SPRING_DARK, alignment=PP_ALIGN.LEFT)

# Q&A slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_text_box(slide, 1, 2.5, 11, 1.5, "Q & A", font_size=64, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.5, 11, 0.8, "有什么想问的？", font_size=28, color=RGBColor(0xCC,0xCC,0xCC), alignment=PP_ALIGN.CENTER)

# Save
output_path = '/Users/lihg/MyOpenCodeProject/springboot-learning/docs/slides.pptx'
prs.save(output_path)
print(f"✅ PPTX saved: {output_path}")
print(f"   Slides: {len(prs.slides)}")
